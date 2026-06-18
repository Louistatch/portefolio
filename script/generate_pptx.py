from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
import json, re

GREEN = RGBColor(0x0D, 0x94, 0x88)
DARK = RGBColor(0x1E, 0x29, 0x3B)
GRAY = RGBColor(0x64, 0x74, 0x8B)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT_BG = RGBColor(0xF1, 0xF5, 0xF9)
ORANGE = RGBColor(0xD9, 0x77, 0x06)
RED = RGBColor(0xDC, 0x26, 0x26)

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

def add_shape(slide, left, top, w, h, fill=None, border=None):
    shape = slide.shapes.add_shape(1, left, top, w, h)  # rectangle
    shape.line.fill.background()
    if fill:
        shape.fill.solid()
        shape.fill.fore_color.rgb = fill
    else:
        shape.fill.background()
    if border:
        shape.line.color.rgb = border
        shape.line.width = Pt(1.5)
    return shape

def set_text(shape, text, size=14, bold=False, color=DARK, align=PP_ALIGN.LEFT):
    shape.text_frame.word_wrap = True
    p = shape.text_frame.paragraphs[0]
    p.text = text
    p.font.size = Pt(size)
    p.font.bold = bold
    p.font.color.rgb = color
    p.alignment = align
    return p

def add_para(tf, text, size=14, bold=False, color=DARK, space_before=6):
    p = tf.add_paragraph()
    p.text = text
    p.font.size = Pt(size)
    p.font.bold = bold
    p.font.color.rgb = color
    p.space_before = Pt(space_before)
    return p

def title_slide(title, subtitle=""):
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank
    bg = add_shape(slide, 0, 0, prs.slide_width, prs.slide_height, fill=GREEN)
    # Title
    txBox = slide.shapes.add_textbox(Inches(1), Inches(2), Inches(11), Inches(2))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(40)
    p.font.bold = True
    p.font.color.rgb = WHITE
    p.alignment = PP_ALIGN.CENTER
    if subtitle:
        p2 = tf.add_paragraph()
        p2.text = subtitle
        p2.font.size = Pt(20)
        p2.font.color.rgb = RGBColor(0xCC, 0xFB, 0xF1)
        p2.alignment = PP_ALIGN.CENTER
        p2.space_before = Pt(16)
    # Footer
    txF = slide.shapes.add_textbox(Inches(1), Inches(6), Inches(11), Inches(0.6))
    pf = txF.text_frame.paragraphs[0]
    pf.text = "DataMEAL Academy — Formation des Formateurs"
    pf.font.size = Pt(12)
    pf.font.color.rgb = RGBColor(0x99, 0xF6, 0xE4)
    pf.alignment = PP_ALIGN.CENTER
    return slide

def section_slide(month_title, subtitle):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    bg = add_shape(slide, 0, 0, prs.slide_width, prs.slide_height, fill=DARK)
    txBox = slide.shapes.add_textbox(Inches(1), Inches(2.5), Inches(11), Inches(2))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = month_title
    p.font.size = Pt(36)
    p.font.bold = True
    p.font.color.rgb = GREEN
    p.alignment = PP_ALIGN.CENTER
    p2 = tf.add_paragraph()
    p2.text = subtitle
    p2.font.size = Pt(20)
    p2.font.color.rgb = RGBColor(0x94, 0xA3, 0xB8)
    p2.alignment = PP_ALIGN.CENTER
    p2.space_before = Pt(12)
    return slide

def lesson_title_slide(num, title, duration, material):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    # Green left bar
    add_shape(slide, 0, 0, Inches(0.15), prs.slide_height, fill=GREEN)
    # Lesson number badge
    badge = add_shape(slide, Inches(1), Inches(1.5), Inches(1.2), Inches(1.2), fill=GREEN)
    set_text(badge, str(num), size=40, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    badge.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
    # Title
    txBox = slide.shapes.add_textbox(Inches(2.8), Inches(1.3), Inches(9.5), Inches(1.5))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = f"Leçon {num}"
    p.font.size = Pt(16)
    p.font.color.rgb = GREEN
    p.font.bold = True
    p2 = tf.add_paragraph()
    p2.text = title
    p2.font.size = Pt(32)
    p2.font.bold = True
    p2.font.color.rgb = DARK
    p2.space_before = Pt(6)
    # Info bar
    info = slide.shapes.add_textbox(Inches(2.8), Inches(3.2), Inches(9.5), Inches(0.8))
    itf = info.text_frame
    itf.word_wrap = True
    ip = itf.paragraphs[0]
    ip.text = f"Durée : {duration}  |  Matériel : {material}"
    ip.font.size = Pt(14)
    ip.font.color.rgb = GRAY
    return slide

def content_slide(title, bullets, callout_title=None, callout_text=None, callout_variant="tip"):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_shape(slide, 0, 0, Inches(0.08), prs.slide_height, fill=GREEN)
    # Title
    txBox = slide.shapes.add_textbox(Inches(0.6), Inches(0.4), Inches(12), Inches(0.8))
    set_text(txBox, title, size=26, bold=True, color=DARK)
    # Bullets
    top = Inches(1.5)
    if callout_title:
        bw = Inches(7)
    else:
        bw = Inches(11.5)
    txB = slide.shapes.add_textbox(Inches(0.6), top, bw, Inches(5))
    tf = txB.text_frame
    tf.word_wrap = True
    for i, b in enumerate(bullets):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = b
        p.font.size = Pt(16)
        p.font.color.rgb = DARK
        p.space_before = Pt(8)
    # Callout box
    if callout_title and callout_text:
        colors = {"tip": GREEN, "warning": ORANGE, "real": RGBColor(0x25, 0x63, 0xEB)}
        c = colors.get(callout_variant, GREEN)
        box = add_shape(slide, Inches(8.2), Inches(1.5), Inches(4.5), Inches(4.5), fill=LIGHT_BG, border=c)
        box.text_frame.word_wrap = True
        pt = box.text_frame.paragraphs[0]
        pt.text = callout_title
        pt.font.size = Pt(14)
        pt.font.bold = True
        pt.font.color.rgb = c
        pc = box.text_frame.add_paragraph()
        pc.text = callout_text
        pc.font.size = Pt(13)
        pc.font.color.rgb = DARK
        pc.space_before = Pt(10)
    return slide

def table_slide(title, headers, rows):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_shape(slide, 0, 0, Inches(0.08), prs.slide_height, fill=GREEN)
    txBox = slide.shapes.add_textbox(Inches(0.6), Inches(0.4), Inches(12), Inches(0.8))
    set_text(txBox, title, size=26, bold=True, color=DARK)
    cols = len(headers)
    r = len(rows) + 1
    col_w = min(Inches(3), Inches(11.5) / cols)
    table_shape = slide.shapes.add_table(r, cols, Inches(0.6), Inches(1.5), col_w * cols, Inches(0.5) * r)
    table = table_shape.table
    for i, h in enumerate(headers):
        cell = table.cell(0, i)
        cell.text = h
        for p in cell.text_frame.paragraphs:
            p.font.size = Pt(13)
            p.font.bold = True
            p.font.color.rgb = WHITE
        cell.fill.solid()
        cell.fill.fore_color.rgb = GREEN
    for ri, row in enumerate(rows):
        for ci, val in enumerate(row):
            cell = table.cell(ri + 1, ci)
            cell.text = str(val)
            for p in cell.text_frame.paragraphs:
                p.font.size = Pt(12)
                p.font.color.rgb = DARK
            cell.fill.solid()
            cell.fill.fore_color.rgb = WHITE if ri % 2 == 0 else LIGHT_BG
    return slide

def quiz_slide(question, options, answer_idx):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_shape(slide, 0, 0, prs.slide_width, prs.slide_height, fill=RGBColor(0xF0, 0xFD, 0xFA))
    add_shape(slide, 0, 0, Inches(0.08), prs.slide_height, fill=GREEN)
    # Quiz badge
    badge = add_shape(slide, Inches(0.6), Inches(0.5), Inches(1), Inches(0.5), fill=GREEN)
    set_text(badge, "QUIZ", size=16, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    # Question
    txQ = slide.shapes.add_textbox(Inches(0.6), Inches(1.3), Inches(12), Inches(1.2))
    set_text(txQ, question, size=22, bold=True, color=DARK)
    txQ.text_frame.word_wrap = True
    # Options
    for i, opt in enumerate(options):
        letter = chr(65 + i)
        is_answer = (i == answer_idx)
        y = Inches(2.8 + i * 0.9)
        box_color = GREEN if is_answer else WHITE
        text_color = WHITE if is_answer else DARK
        box = add_shape(slide, Inches(1.5), y, Inches(10), Inches(0.7), fill=box_color, border=GREEN if not is_answer else None)
        set_text(box, f"  {letter}.  {opt}", size=16, bold=is_answer, color=text_color)
    return slide

# ══════════════════════════════════════════════════
# BUILD THE PRESENTATION
# ══════════════════════════════════════════════════

# --- COVER ---
title_slide(
    "Formation des Formateurs\nGestion Financière Paysanne",
    "Parcours andragogique de 3 mois — 12 sessions\nDataMEAL Academy | TOF-FIN-01"
)

# --- OVERVIEW ---
table_slide("Programme — 12 sessions sur 3 mois",
    ["#", "Période", "Thème", "Activité pratique"],
    [
        ["1", "Mois 1 - S1", "Posture du formateur rural", "Cercle de parole"],
        ["2", "Mois 1 - S2", "D'où vient et où va l'argent ?", "Cailloux entrants/sortants"],
        ["3", "Mois 1 - S3", "Budget familial en graines", "12 tas + graines maïs/haricots"],
        ["4", "Mois 1 - S4", "Épargne, tontine et AVEC", "Simulation tontine en cercle"],
        ["5", "Mois 2 - S5", "Coût de production agricole", "Calcul collectif au tableau"],
        ["6", "Mois 2 - S6", "Planifier la campagne", "Corde des saisons"],
        ["7", "Mois 2 - S7", "Le crédit : opportunité ou piège ?", "Jeu de rôle 3 prêteurs"],
        ["8", "Mois 2 - S8", "Vendre au bon moment", "Courbe des prix au sol"],
        ["9", "Mois 3 - S9", "Gérer les risques financiers", "Méthode des 3 enveloppes"],
        ["10", "Mois 3 - S10", "Structurer une session APCA", "Préparer une fiche de session"],
        ["11", "Mois 3 - S11", "Pratique supervisée", "Chaque stagiaire anime 20 min"],
        ["12", "Mois 3 - S12", "Plan d'action communautaire", "Engagement collectif"],
    ]
)

# --- APPROCHE ---
content_slide("Approche andragogique : 5 principes",
    [
        "1. L'adulte a besoin de savoir POURQUOI → Commencer par un problème concret",
        "2. L'adulte apporte son expérience → Faire parler les anciens d'abord",
        "3. L'adulte veut décider lui-même → Proposer, ne jamais imposer",
        "4. L'adulte apprend pour résoudre un problème immédiat → Chaque session = 1 problème réel",
        "5. L'adulte est motivé par des résultats concrets → Gains visibles dès la 1ère session",
    ],
    callout_title="Règle d'or",
    callout_text="Le formateur parle 20% du temps maximum. Les 80% restants, ce sont les participants qui parlent, réfléchissent, calculent et échangent.",
    callout_variant="tip"
)

# Cycle APCA
content_slide("Le cycle APCA — Canevas de chaque session",
    [
        "A — ACCROCHE (15 min) : Un problème vécu, une histoire, une question provocante",
        "",
        "P — PRATIQUE (60 min) : Activité participative (cailloux, jeu de rôle, calcul collectif)",
        "",
        "C — CONCEPTUALISATION (20 min) : Le GROUPE formule la leçon — pas le formateur",
        "",
        "A — ACTION (15 min) : Engagement concret à appliquer avant la prochaine session",
        "",
        "Durée totale : 2h maximum. Au-delà, les adultes en milieu rural décrochent.",
    ],
    callout_title="Pourquoi APCA ?",
    callout_text="Les adultes retiennent 20% de ce qu'ils entendent, 50% de ce qu'ils voient, mais 80% de ce qu'ils FONT. Le cycle APCA met la PRATIQUE au centre.",
    callout_variant="real"
)

# ═══ MOIS 1 ═══
section_slide("MOIS 1", "Comprendre l'argent du paysan\nSemaines 1 à 4")

# L1
lesson_title_slide(1, "Posture du formateur-animateur\nen milieu rural", "2h", "cercle de chaises, images parlantes")
content_slide("Vous n'êtes pas un professeur, vous êtes un facilitateur",
    [
        "Les paysans ont 20, 30, 40 ans d'expérience. Ils savent déjà beaucoup.",
        "Votre rôle : faciliter la prise de conscience, PAS enseigner.",
        "",
        "Boîte à outils du facilitateur :",
        "  • L'arbre à problèmes — dessiner causes/conséquences au sol",
        "  • Le cercle de parole — chacun parle à son tour",
        "  • Les images parlantes — dessins sans texte, interprétés par le groupe",
        "  • Le jeu de rôle — simuler des situations financières réelles",
        "  • Les cailloux/graines — compter, représenter, répartir",
    ],
    callout_title="Le piège du débutant",
    callout_text="Un formateur arrive avec des PowerPoint et du jargon comptable. Les paysans s'ennuient et ne reviennent pas. Ce qui marche : s'asseoir sous l'arbre à palabres, écouter d'abord.",
    callout_variant="warning"
)
quiz_slide("Un formateur rural efficace consacre quelle proportion du temps à écouter et faire parler le groupe ?",
    ["20% écoute, 80% cours magistral", "50% écoute, 50% cours", "80% écoute et facilitation, 20% apport", "100% cours théorique"], 2)

# L2
lesson_title_slide(2, "D'où vient et où va l'argent\ndu paysan ?", "2h", "sol sableux, cailloux 2 couleurs, bâton")
content_slide("Activité terrain : la carte des flux d'argent",
    [
        "1. Dessinez un grand cercle au sol = la famille / le ménage",
        "2. « D'où vient l'argent ? » → cailloux verts ENTRANTS",
        "3. « Où part l'argent ? » → cailloux rouges SORTANTS",
        "4. Le groupe débat et complète",
        "",
        "Les 4 sources de revenus typiques :",
        "  • Vente de récolte (1-2x/an, gros montant)",
        "  • Petit commerce (quotidien, petits montants)",
        "  • Travail journalier (saisonnier, variable)",
        "  • Transferts diaspora (irrégulier)",
    ],
    callout_title="Le problème fondamental",
    callout_text="Les REVENUS sont saisonniers (1-2 pics par an) mais les DÉPENSES sont permanentes (nourriture, santé, école tous les mois). C'est le déséquilibre que l'activité fait apparaître.",
    callout_variant="real"
)
table_slide("Le déséquilibre saisonnier — revenus vs dépenses",
    ["Mois", "Revenus (FCFA)", "Dépenses (FCFA)", "Solde"],
    [
        ["Jan–Mar", "45 000", "110 000", "-65 000  DÉFICIT"],
        ["Avr–Jun", "20 000", "145 000", "-125 000  DÉFICIT"],
        ["Jul–Sep", "20 000", "125 000", "-105 000  DÉFICIT"],
        ["Oct–Déc", "330 000", "125 000", "+205 000"],
        ["TOTAL", "415 000", "505 000", "Le paysan traverse 9 mois de déficit !"],
    ]
)
quiz_slide("Quel est le problème financier fondamental que l'activité des cailloux fait apparaître ?",
    ["Il ne gagne pas assez", "Ses revenus sont saisonniers mais ses dépenses sont permanentes",
     "Il dépense trop en cérémonies", "Il ne sait pas compter"], 1)

# L3
lesson_title_slide(3, "Construire un budget familial\navec les paysans", "3h", "graines (maïs, haricots), 12 cuvettes, feuille A3")
content_slide("L'outil central : le budget en tas de graines",
    [
        "1. Disposez 12 tas de sable = janvier à décembre",
        "2. 100 graines de maïs = « tout l'argent gagné en un an »",
        "3. « Répartissez les graines dans les mois où vous recevez de l'argent »",
        "4. 100 graines de haricots = les dépenses par mois",
        "5. Le paysan VOIT les mois où les haricots dépassent le maïs",
        "",
        "Ensuite, pour ceux qui le souhaitent → tableau papier simplifié :",
        "  Mois | Ce qui rentre | Ce qui sort | Il reste",
        "",
        "Les montants viennent DU PAYSAN, pas de vous.",
    ],
    callout_title="Témoignage terrain (Mali)",
    callout_text="« Quand Mamadou a vu que pendant 7 mois sur 12 il avait plus de haricots que de maïs, il a dit : Ah, c'est pour ça qu'on souffre de mars à septembre ! Il le vivait chaque année mais ne l'avait jamais VU. »",
    callout_variant="real"
)
quiz_slide("Pourquoi utiliser des graines plutôt qu'un tableau écrit pour le budget ?",
    ["C'est moins cher", "Ça permet aux non-alphabétisés de participer et de VOIR le déséquilibre",
     "Les graines sont plus jolies", "C'est obligatoire"], 1)

# L4
lesson_title_slide(4, "L'épargne paysanne :\ntontines, AVEC et stratégies locales", "2h30", "dessins (grenier, tontine, banque, téléphone), graines")
content_slide("Les paysans épargnent DÉJÀ — mais comment ?",
    [
        "Tour de cercle : « Où gardez-vous votre argent ? »",
        "  • Sous le matelas / pot enterré",
        "  • Grenier (stock de grain = épargne en nature !)",
        "  • Chez un parent de confiance",
        "  • Tontine du village",
        "  • Mobile money",
        "",
        "La tontine : 10 membres × 5 000 FCFA/semaine",
        "  → 50 000 FCFA pour UN membre chaque semaine",
        "  → Au bout de 10 semaines, tout le monde a reçu une fois",
        "",
        "Simulation : 10 volontaires en cercle avec des graines",
        "Question : « À quoi sert la tontine si on reçoit ce qu'on donne ? »",
    ],
    callout_title="Le grenier = un compte épargne",
    callout_text="Stocker 5 sacs de mil à 250 FCFA/kg en novembre → revente à 400 FCFA/kg en juillet. Rendement de 60% ! Les paysans le font depuis des générations. VALORISEZ cette intelligence.",
    callout_variant="tip"
)
table_slide("Comparaison des formes d'épargne",
    ["Type", "Avantage", "Risque", "Pour qui"],
    [
        ["Sous le matelas", "Immédiatement disponible", "Vol, tentation", "Dépannage court terme"],
        ["Grenier (stock grain)", "+60% de valeur en 6 mois", "Insectes, incendie", "Paysans avec stockage sec"],
        ["Tontine simple", "Force l'épargne, gros montant", "Défaillance d'un membre", "Tous"],
        ["AVEC / VSLA", "Épargne + crédit + fonds social", "Nécessite formation", "Groupes organisés"],
        ["Mobile money", "Sécurisé, accessible", "Frais, réseau nécessaire", "Zones couvertes"],
    ]
)
quiz_slide("Quel est le principal avantage de la tontine pour un paysan ?",
    ["Elle rapporte des intérêts", "Elle force l'épargne régulière et donne un gros montant d'un coup",
     "Elle remplace la banque", "Elle est gratuite"], 1)

# ═══ MOIS 2 ═══
section_slide("MOIS 2", "Gérer et planifier\nSemaines 5 à 8")

# L5
lesson_title_slide(5, "Calculer le coût de production\navec les paysans", "3h", "grande feuille, graines, dessins (champ, engrais, houe)")
content_slide("Le paysan connaît-il le vrai coût de sa production ?",
    [
        "Animation : « Combien vous coûte un sac de maïs ? »",
        "1. Dessinez un grand champ au tableau/sol",
        "2. « Pour cultiver 1 ha de maïs, qu'est-ce qu'il vous faut ? »",
        "3. Le groupe liste tout : semences, engrais, labour, sarclage, récolte...",
        "4. Pour CHAQUE élément : « Combien ça coûte ? »",
        "5. Additionnez ensemble. La surprise est souvent énorme.",
        "",
        "Le coût invisible : la main-d'œuvre familiale",
        "« Si votre femme allait vendre au marché au lieu de sarcler,",
        "  combien gagnerait-elle ? » → C'est le coût d'opportunité.",
    ],
    callout_title="Le choc du calcul",
    callout_text="Au nord-Cameroun, un paysan vendait son sac à 8 000 FCFA en pensant gagner de l'argent. Le coût réel : 9 500 FCFA/sac. Il travaillait à PERTE depuis des années sans le savoir.",
    callout_variant="warning"
)
table_slide("Coût de production — 1 hectare de maïs (zone soudanienne)",
    ["Poste", "Montant (FCFA)"],
    [
        ["Semences (20 kg)", "15 000"],
        ["Engrais NPK (2 sacs)", "40 000"],
        ["Engrais urée (1 sac)", "22 000"],
        ["Labour (tracteur/attelage)", "25 000"],
        ["Semis (main-d'œuvre)", "10 000"],
        ["1er sarclage", "15 000"],
        ["2e sarclage", "12 000"],
        ["Récolte", "20 000"],
        ["Transport + stockage", "13 000"],
        ["Produit conservation", "3 000"],
        ["TOTAL", "175 000"],
        ["Coût par sac (20 sacs/ha)", "8 750 FCFA → prix plancher de vente"],
    ]
)
quiz_slide("Un paysan vend son sac à 8 000 FCFA. Le coût total est de 9 500 FCFA/sac. Quelle est la situation ?",
    ["Bénéfice de 1 500 FCFA", "Perte de 1 500 FCFA par sac",
     "Équilibre", "Impossible à dire"], 1)

# L6
lesson_title_slide(6, "Planifier financièrement\nla campagne agricole", "3h", "corde 3m, 12 nœuds, cartons illustrés des activités")
content_slide("Le calendrier cultural = le calendrier financier",
    [
        "Activité : La corde des saisons",
        "1. Corde avec 12 nœuds (1/mois) entre deux piquets",
        "2. Placez les cartons : labour, semis, sarclage, récolte, vente",
        "3. Sous chaque carton : graines = argent nécessaire",
        "4. Question clé : « En avril, quand il faut acheter intrants,",
        "   est-ce que vous avez l'argent ? »",
        "",
        "La stratégie d'épargne anticipée :",
        "  Besoins en intrants (avril) : ~80 000 FCFA",
        "  Mois pour épargner (nov→mars) : 5 mois",
        "  Épargne nécessaire : 80 000 ÷ 5 = 16 000 FCFA/mois",
        "  → Concret, atteignable, brise le cycle d'endettement",
    ],
    callout_title="Le cercle vicieux",
    callout_text="En avril, plus d'argent → emprunt au commerçant → remboursement en nature à la récolte (prix le plus bas) → perte de 30-50% de la valeur. Chaque année, le même piège. La planification le brise.",
    callout_variant="warning"
)
quiz_slide("Pourquoi le paysan s'endette-t-il souvent en avril ?",
    ["Parce qu'il est paresseux",
     "Parce qu'il n'a pas planifié l'épargne pour couvrir les intrants",
     "Parce que les intrants sont trop chers",
     "Parce que la banque le force"], 1)

# L7
lesson_title_slide(7, "Le crédit agricole :\nopportunité ou piège ?", "2h30", "3 chaises, cartons avec offres de crédit")
content_slide("Jeu de rôle : le paysan face au crédit",
    [
        "Scénario : Aminata a besoin de 100 000 FCFA pour les intrants",
        "",
        "3 prêteurs se présentent :",
        "  A. Le commerçant : intrants maintenant, 5 sacs à la récolte",
        "     → Valeur remboursée : 150 000 FCFA — Coût : 50 000 (50%)",
        "  B. La microfinance : 2%/mois sur 8 mois",
        "     → Remboursement : 119 000 FCFA — Coût : 19 000 (19%)",
        "  C. L'AVEC/VSLA : 5% sur 4 mois, flexible",
        "     → Remboursement : 105 000 FCFA — Coût : 5 000 (5%)",
        "",
        "Les 3 règles d'or du crédit :",
        "  1. Emprunter UNIQUEMENT pour un investissement productif",
        "  2. Toujours calculer le coût TOTAL du remboursement",
        "  3. Préférer l'épargne au crédit quand c'est possible",
    ],
    callout_title="Pas de jugement",
    callout_text="Si un participant dit 'moi j'emprunte toujours au commerçant', ne le critiquez pas. Demandez 'combien tu lui rembourses ?' et laissez le groupe calculer. Le jugement vient du groupe, pas du formateur.",
    callout_variant="tip"
)
quiz_slide("Un commerçant prête 100 000 FCFA en intrants et demande 5 sacs à 30 000 FCFA. Quel est le coût du crédit ?",
    ["0 FCFA, c'est un échange", "50 000 FCFA (150 000 - 100 000)",
     "30 000 FCFA", "100 000 FCFA"], 1)

# L8
lesson_title_slide(8, "Stratégies de commercialisation :\nvendre au bon moment", "2h30", "cailloux pour courbe des prix, 12 mois au sol")
content_slide("Le paradoxe : il vend quand le prix est le plus bas",
    [
        "Activité : La courbe des prix au sol",
        "1. Ligne de 3 mètres = 12 mois (janvier à décembre)",
        "2. « En novembre, un sac coûte combien ? Et en juillet ? »",
        "3. Empilez des cailloux = le prix de chaque mois",
        "4. Le groupe VOIT : prix bas à la récolte, prix haut en soudure",
        "",
        "Les 4 stratégies de commercialisation :",
        "  • Vente groupée — vendre ensemble = meilleur prix",
        "  • Stockage individuel — garder et vendre en mai-juillet",
        "  • Warrantage — stocker + crédit 70% + vendre plus tard",
        "  • Transformation — farine, huile = valeur ajoutée",
    ],
    callout_title="Le warrantage (Niger/PAM)",
    callout_text="Les paysans stockent leur grain dans un entrepôt certifié, reçoivent 70% en crédit immédiat, et vendent 4-6 mois plus tard. Résultat : +40% de revenus en moyenne.",
    callout_variant="real"
)
table_slide("Gain du stockage — 20 sacs de maïs",
    ["Scénario", "Prix/sac", "Total", "Différence"],
    [
        ["Vente immédiate (novembre)", "8 000 FCFA", "160 000 FCFA", "—"],
        ["Vente différée (mai)", "17 000 FCFA", "340 000 FCFA", "+180 000 FCFA"],
        ["Coût de stockage", "—", "-15 000 FCFA", "—"],
        ["GAIN NET du stockage", "—", "—", "+165 000 FCFA (+103%)"],
    ]
)
quiz_slide("Un paysan vend 20 sacs à 8 000 FCFA au lieu de 17 000 FCFA. Combien perd-il ?",
    ["Rien", "90 000 FCFA", "180 000 FCFA (9 000 × 20 sacs)", "20 000 FCFA"], 2)

# ═══ MOIS 3 ═══
section_slide("MOIS 3", "Transmettre et pérenniser\nSemaines 9 à 12")

# L9
lesson_title_slide(9, "Gérer les risques :\nmaladies, sécheresse, vol", "2h30", "cartes illustrées des risques, 3 boîtes")
content_slide("Les risques qui ruinent les paysans",
    [
        "Activité : le classement des risques",
        "  Cartes illustrées : sécheresse, maladie, vol, chute des prix...",
        "  « Quel est le plus fréquent ? Le plus grave ? »",
        "",
        "3 stratégies de gestion :",
        "  • PRÉVENIR — diversifier les cultures, variétés résistantes",
        "  • ATTÉNUER — fonds d'urgence (épargne séparée)",
        "  • TRANSFÉRER — assurance, mutuelle santé, fonds social AVEC",
        "",
        "Priorité n°1 : le fonds d'urgence",
        "  = 2-3 mois de dépenses alimentaires, SÉPARÉ de l'épargne intrants",
        "",
        "Méthode des 3 enveloppes/boîtes : INTRANTS | URGENCE | FAMILLE",
    ],
    callout_title="L'histoire de Moussa",
    callout_text="Moussa avait 200 000 FCFA pour un bœuf. Son fils tombe malade. Toute l'épargne part en soins. Plus de bœuf, campagne ratée. Sans fonds d'urgence séparé, UN événement efface des ANNÉES d'efforts.",
    callout_variant="warning"
)
quiz_slide("Pourquoi le fonds d'urgence doit-il être SÉPARÉ de l'épargne pour les intrants ?",
    ["Pour compliquer les choses",
     "Pour éviter qu'un imprévu détruise la campagne en entamant l'épargne productive",
     "Pour gagner plus d'intérêts",
     "Ce n'est pas nécessaire"], 1)

# L10
lesson_title_slide(10, "Structurer et animer une session\nde formation au village", "2h", "fiches APCA vierges, chronomètre")
content_slide("Le plan de session APCA : votre canevas",
    [
        "Chaque session = 4 phases, 2h maximum :",
        "",
        "  ACCROCHE (15 min)",
        "  → Problème vécu, histoire, question provocante",
        "",
        "  PRATIQUE (60 min)",
        "  → Activité participative : cailloux, jeu de rôle, calcul collectif",
        "",
        "  CONCEPTUALISATION (20 min)",
        "  → Le GROUPE formule la leçon (pas le formateur !)",
        "",
        "  ACTION (15 min)",
        "  → Engagement concret avant la prochaine session",
    ],
    callout_title="La fiche de session",
    callout_text="Pour CHAQUE session, préparez une fiche A4 recto-verso : objectif, matériel, déroulement APCA minute par minute, questions de facilitation. C'est l'outil de démultiplication que vous donnez aux formateurs que VOUS formez.",
    callout_variant="tip"
)
table_slide("Exemple : session sur le budget familial",
    ["Phase", "Durée", "Le formateur fait", "Les participants font"],
    [
        ["Accroche", "15 min", "Raconte l'histoire de Mamadou", "Réagissent, s'identifient"],
        ["Pratique", "60 min", "Distribue les graines, guide", "Construisent LEUR budget"],
        ["Conceptualisation", "20 min", "« Qu'avez-vous découvert ? »", "Formulent la leçon eux-mêmes"],
        ["Action", "15 min", "« Que comptez-vous faire ? »", "S'engagent concrètement"],
    ]
)
quiz_slide("Dans le cycle APCA, à quel moment le formateur parle-t-il le plus ?",
    ["Pendant l'Accroche", "Pendant la Pratique",
     "Nulle part — 20% max, le groupe fait 80%",
     "Pendant la Conceptualisation"], 2)

# L11
lesson_title_slide(11, "Pratique supervisée :\nanimez votre première session", "3h", "tout le matériel du thème choisi")
content_slide("C'est à vous de jouer !",
    [
        "Consignes :",
        "  1. Choisissez UN thème (sessions 1 à 9)",
        "  2. Préparez votre fiche APCA",
        "  3. Animez 20 minutes devant vos co-stagiaires",
        "  4. Feedback constructif du groupe + superviseur",
        "",
        "Grille d'observation (5 critères, note 1 à 3) :",
        "  • Pose des questions avant d'expliquer",
        "  • Utilise du matériel concret",
        "  • Langage simple, pas de jargon",
        "  • Respect du cycle APCA et du temps",
        "  • Bienveillance, valorise les savoirs du groupe",
        "",
        "Feedback : méthode sandwich (positif → amélioration → positif)",
    ],
    callout_title="L'erreur la plus fréquente",
    callout_text="90% des stagiaires font la même erreur : ils parlent trop. C'est normal — on a tous été formés par le modèle scolaire. L'andragogie c'est le contraire. Il faut PRATIQUER pour désapprendre.",
    callout_variant="warning"
)
quiz_slide("Lors du feedback à un stagiaire, par quoi commencez-vous ?",
    ["Par les erreurs à corriger", "Par ce qu'il a bien fait (méthode sandwich)",
     "Par une note chiffrée", "Par un silence"], 1)

# L12 — CAPSTONE
lesson_title_slide(12, "Capstone : votre plan d'action\ncommunautaire", "3h", "fiches plan d'action, tableau de suivi")
content_slide("Bilan et engagement — que va changer ce parcours ?",
    [
        "Ce que vous maîtrisez maintenant :",
        "  ✅ Posture de facilitateur (pas de professeur)",
        "  ✅ Diagnostic financier participatif (cailloux, graines)",
        "  ✅ Budget familial et épargne (tontine, AVEC)",
        "  ✅ Coût de production et planification de campagne",
        "  ✅ Crédit, commercialisation, gestion des risques",
        "  ✅ Structurer et animer une session APCA",
        "",
        "Votre plan d'action :",
        "  1. OÙ — dans quel(s) village(s) ?",
        "  2. QUAND — calendrier des 12 sessions",
        "  3. QUI — combien de groupes et de participants ?",
        "  4. COMMENT — matériel nécessaire",
        "  5. SUIVI — mesurer le changement (avant/après)",
    ],
    callout_title="Impact projeté",
    callout_text="1 formateur × 4 groupes × 25 paysans = 100 ménages. Si chaque paysan en forme 2 : 300 ménages touchés = 2 100 personnes. +35% épargne, -40% crédit commerçant après 1 an.",
    callout_variant="real"
)
table_slide("Tableau de suivi d'impact",
    ["Indicateur", "Avant formation", "Après 6 mois", "Objectif"],
    [
        ["% ménages avec un budget", "À mesurer", "À mesurer", "60%"],
        ["Épargne moyenne mensuelle", "À mesurer", "À mesurer", "+30%"],
        ["% utilisant le crédit commerçant", "À mesurer", "À mesurer", "-50%"],
        ["Nombre de tontines/AVEC actives", "À compter", "À compter", "+3"],
        ["Paysans formés (direct + indirect)", "0", "À compter", "300"],
    ]
)

# --- CLOSING ---
title_slide(
    "Vous êtes maintenant\nFormateur en Gestion Financière Paysanne",
    "Allez sur le terrain. Écoutez les paysans.\nAccompagnez-les avec des cailloux, des graines et beaucoup d'écoute.\n\nBon terrain !"
)

# Save
output = "/home/user/portefolio/TOF-FIN-01_Gestion_Financiere_Paysanne.pptx"
prs.save(output)
print(f"PPTX saved: {output}")
